import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def extract_content_from_pdf(file_bytes):
    """Extracts text and unique, significant image count from a PDF."""
    text = ""
    unique_images = set()
    
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
            # get_images returns a list of tuples: (xref, smask, width, height, bpc, colorspace, ...)
            for img in page.get_images(full=True):
                xref = img[0]
                width = img[2]
                height = img[3]
                
                # Filter: Only count "real" images (ignore icons/lines < 100x100 pixels)
                # and deduplicate by xref (so logos on every page only count once)
                if width > 100 and height > 100:
                    unique_images.add(xref)
                    
    return text, len(unique_images)

def extract_content_from_pptx(file_bytes):
    """Extracts text and unique image count from a PPTX."""
    text = ""
    image_blobs = set() # Use blobs to deduplicate images in PPTX
    
    prs = Presentation(io.BytesIO(file_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Deduplicate by hashing the image binary data
                img_blob = shape.image.blob
                if shape.width.inches > 0.5 and shape.height.inches > 0.5: # Simple size filter
                    image_blobs.add(hash(img_blob))
            
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if s.width.inches > 0.5 and s.height.inches > 0.5:
                            image_blobs.add(hash(s.image.blob))
                            
    return text, len(image_blobs)

def extract_content(file_name, file_bytes):
    """Dispatches extraction based on file extension. Returns (text, image_count)."""
    ext = file_name.split('.')[-1].lower()
    
    # Text-based formats (including code and structured data)
    text_extensions = [
        'txt', 'md', 'html', 'htm', 'json', 'csv', 'xml', 'yaml', 'yml',
        'py', 'js', 'ts', 'jsx', 'tsx', 'c', 'cpp', 'h', 'hpp', 'java', 'go', 'rb', 'php', 'sql'
    ]
    
    if ext in text_extensions:
        return file_bytes.decode('utf-8', errors='ignore'), 0
    elif ext == 'pdf':
        return extract_content_from_pdf(file_bytes)
    elif ext == 'pptx':
        return extract_content_from_pptx(file_bytes)
    elif ext in ['png', 'jpg', 'jpeg', 'webp']:
        return "", 1
    else:
        return "", 0
